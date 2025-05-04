import os
import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from pptx.util import Inches

# Função para abrir a janela de seleção de arquivos
def selecionar_imagens():
    # Cria a janela de seleção de arquivos
    root = tk.Tk()
    root.withdraw()  # Não mostrar a janela principal
    arquivos = filedialog.askopenfilenames(
        title="Selecione as imagens para montar uma apresentação de slides", 
        filetypes=[("Arquivos de Imagem", "*.jpg;*.jpeg;*.png")]
    )
    return arquivos

# Função para criar um PowerPoint com as imagens selecionadas
def criar_powerpoint(imagens):
    print("Criando apresentação...")
    # Verifica se há imagens selecionadas
    if not imagens:
        print("Nenhuma imagem selecionada.")
        return
    # Cria uma apresentação
    ppt = Presentation()

    # Define o tamanho do slide para 16:9 (13.33 x 7.5 polegadas)
    ppt.slide_width = Inches(13.33)
    ppt.slide_height = Inches(7.5)

    # Definir as dimensões do slide
    image_width = ppt.slide_width
    image_height = int(ppt.slide_width * 9 / 16) # Proporção 16:9

    # Adiciona cada imagem como um slide
    for imagem in imagens:
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Layout em branco
        slide.shapes.add_picture(imagem, Inches(0), Inches(0), width=image_width, height=image_height)
        print(f"Adicionando imagem: {imagem}")

    # Salva o PowerPoint no mesmo diretório onde o script está
    diretorio_atual = os.path.dirname(os.path.abspath(__file__))
    nome_arquivo = os.path.join(diretorio_atual, "apresentacao_imagens.pptx")
    ppt.save(nome_arquivo)
    print(f"PowerPoint salvo como: {nome_arquivo}")

# Função principal
def main():
    imagens = selecionar_imagens()  # Seleciona as imagens
    if imagens:  # Verifica se alguma imagem foi selecionada
        print(f"{len(imagens)} imagem(s) selecionada(s).")
        criar_powerpoint(imagens)  # Cria o PowerPoint com as imagens
    else:
        print("Nenhuma imagem foi selecionada.")

# Executa o script
if __name__ == "__main__":
    main()